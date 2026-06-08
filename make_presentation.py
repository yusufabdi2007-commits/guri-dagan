from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Colours ──────────────────────────────────────────────────────────────────
C_BG        = RGBColor(0x0F, 0x17, 0x2A)   # dark navy
C_ACCENT    = RGBColor(0x6C, 0x63, 0xFF)   # purple
C_ACCENT2   = RGBColor(0x10, 0xB9, 0x81)   # emerald
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xC7, 0xD2, 0xFE)   # lavender text
C_CARD      = RGBColor(0x1E, 0x29, 0x45)   # card bg
C_YELLOW    = RGBColor(0xFB, 0xBF, 0x24)
C_ROSE      = RGBColor(0xF4, 0x3F, 0x5E)
C_SKY       = RGBColor(0x38, 0xBD, 0xF8)

def set_bg(slide, color=C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, radius=False):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, text, l, t, w, h,
                 size=18, bold=False, color=C_WHITE,
                 align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_bullet_box(slide, items, l, t, w, h, size=16, color=C_WHITE, dot_color=C_ACCENT):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # dot
        dot = p.add_run()
        dot.text = "● "
        dot.font.size = Pt(size - 2)
        dot.font.color.rgb = dot_color
        # text
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_after = Pt(6)

def add_tag(slide, text, l, t, color, text_color=C_WHITE, size=13):
    w = len(text) * 0.11 + 0.3
    rect = add_rect(slide, l, t, w, 0.35, color)
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = text_color

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

# big gradient-feel bar
add_rect(slide, 0, 0, 13.33, 0.08, C_ACCENT)

# circle decoration
circ = slide.shapes.add_shape(9, Inches(9.5), Inches(0.5), Inches(5), Inches(5))
circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x1E,0x29,0x45)
circ.line.fill.background()

circ2 = slide.shapes.add_shape(9, Inches(10.5), Inches(1.5), Inches(3), Inches(3))
circ2.fill.solid(); circ2.fill.fore_color.rgb = RGBColor(0x6C,0x63,0xFF)
circ2.fill.fore_color.theme_color  # ignore
circ2.line.fill.background()

add_text_box(slide, "Guri Dagan", 0.6, 1.5, 8, 1.4, size=60, bold=True, color=C_WHITE)
add_text_box(slide, "Parenting Coach Platform", 0.65, 2.9, 8, 0.7, size=26, color=C_LIGHT)
add_text_box(slide, "A simple guide to every screen and button", 0.65, 3.7, 8, 0.6, size=18, color=RGBColor(0x94,0xA3,0xB8))
add_text_box(slide, "guri-dagan.vercel.app", 0.65, 6.5, 5, 0.5, size=14, color=C_ACCENT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — How to open the app
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.08, C_ACCENT2)
add_text_box(slide, "How to Open the App", 0.6, 0.3, 12, 0.8, size=36, bold=True, color=C_WHITE)
add_text_box(slide, "Install it on your iPhone — no App Store needed", 0.65, 1.0, 12, 0.5, size=18, color=C_LIGHT)

steps = [
    ("1", "Open Safari on your iPhone"),
    ("2", 'Type  guri-dagan.vercel.app  in the address bar and tap Go'),
    ("3", "Log in with your email and password"),
    ("4", "Tap the Share button (box with arrow) at the bottom of Safari"),
    ("5", 'Tap "Add to Home Screen"  then tap  "Add"'),
    ("6", "Done! The app icon will appear on your home screen"),
]
for i, (num, txt) in enumerate(steps):
    y = 1.7 + i * 0.75
    add_rect(slide, 0.6, y, 0.5, 0.5, C_ACCENT)
    add_text_box(slide, num, 0.6, y, 0.5, 0.5, size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, txt, 1.3, y + 0.05, 11, 0.45, size=17, color=C_WHITE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Bottom Navigation
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.08, C_YELLOW)
add_text_box(slide, "Bottom Navigation Bar", 0.6, 0.3, 12, 0.8, size=36, bold=True, color=C_WHITE)
add_text_box(slide, "The 5 buttons always visible at the bottom of every screen", 0.65, 1.0, 12, 0.5, size=18, color=C_LIGHT)

nav_items = [
    ("Today",     C_ACCENT2,  "Your daily task — one post to focus on today"),
    ("Ideas",     C_ACCENT,   "Save content ideas — topics to record videos about"),
    ("Videos",    C_SKY,      "Your video library — all recorded and posted videos"),
    ("Dashboard", C_YELLOW,   "Overview — stats, streak, and quick actions"),
    ("More",      C_ROSE,     "Access all other pages: Calendar, CRM, Analytics, etc."),
]
for i, (name, col, desc) in enumerate(nav_items):
    x = 0.6 + i * 2.4
    add_rect(slide, x, 1.8, 2.1, 0.55, col)
    add_text_box(slide, name, x, 1.8, 2.1, 0.55, size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, x, 2.5, 2.1, 1.2, size=13, color=C_LIGHT)

add_text_box(slide, "Tip: Tap \"More\" to find ALL the extra pages like Leads, Analytics, Weekly Report, and more.",
             0.6, 4.2, 12, 0.7, size=15, color=C_YELLOW)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Today Page
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.08, C_ACCENT2)
add_text_box(slide, "Today  —  Your Daily Focus", 0.6, 0.3, 12, 0.8, size=36, bold=True, color=C_WHITE)
add_text_box(slide, "The first screen that opens when you launch the app", 0.65, 1.0, 12, 0.5, size=18, color=C_LIGHT)

add_text_box(slide, "What you see:", 0.6, 1.6, 6, 0.4, size=16, bold=True, color=C_ACCENT2)
add_bullet_box(slide, [
    "Today's post — the video or idea scheduled for today",
    "\"Mark as Posted\" button — tap this when you've uploaded the video",
    "Posting checklist — steps to complete before posting (hook, caption, hashtags)",
    "Low energy mode — tap if you're tired; shows a smaller, easier task",
    "AI suggestion — one personalised action for the day",
], 0.6, 2.0, 6.2, 3.5, size=15, dot_color=C_ACCENT2)

add_text_box(slide, "What to do:", 6.9, 1.6, 6, 0.4, size=16, bold=True, color=C_YELLOW)
add_bullet_box(slide, [
    "Open the app every morning and check Today first",
    "Complete the checklist before posting",
    "Tap 'Mark as Posted' after uploading — this builds your streak",
    "If busy, tap 'Low Energy Mode' for a simpler task",
], 6.9, 2.0, 6.0, 3.0, size=15, dot_color=C_YELLOW)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Dashboard
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.08, C_YELLOW)
add_text_box(slide, "Dashboard  —  Your Overview", 0.6, 0.3, 12, 0.8, size=36, bold=True, color=C_WHITE)
add_text_box(slide, "See everything at a glance", 0.65, 1.0, 12, 0.5, size=18, color=C_LIGHT)

cards = [
    ("Streak Ring",     C_YELLOW,  "Shows how many days in a row you've posted. Keep it going!"),
    ("Posted Today",    C_ACCENT2, "Green tick when you've posted today. Stays grey if not yet."),
    ("Weekly Goal",     C_ACCENT,  "How many videos you've posted this week vs your goal."),
    ("Quick Actions",   C_SKY,     "Shortcuts to Ideas, Videos, Calendar, and Generator."),
    ("AI Coach Card",   C_ROSE,    "A daily message from the AI coach with a tip or motivation."),
    ("Momentum",        C_YELLOW,  "Today's AI suggestion — one action to keep your momentum."),
]
for i, (title, col, desc) in enumerate(cards):
    col_x = 0.6 + (i % 3) * 4.1
    row_y = 1.7 + (i // 3) * 2.2
    add_rect(slide, col_x, row_y, 3.8, 0.45, col)
    add_text_box(slide, title, col_x, row_y, 3.8, 0.45, size=15, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, col_x, row_y + 0.5, 3.8, 1.5, size=13, color=C_LIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Ideas Page
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.08, C_ACCENT)
add_text_box(slide, "Ideas  —  Your Content Bank", 0.6, 0.3, 12, 0.8, size=36, bold=True, color=C_WHITE)
add_text_box(slide, "Save every topic or idea so you never run out of things to record", 0.65, 1.0, 12, 0.5, size=18, color=C_LIGHT)

add_text_box(slide, "Buttons on this page:", 0.6, 1.6, 6, 0.4, size=16, bold=True, color=C_ACCENT)
add_bullet_box(slide, [
    "Microphone button (bottom right) — speak your idea out loud, AI writes it for you",
    "Plus (+) button — type an idea manually",
    "Generate button — AI creates new content ideas based on your topic",
    "Filter tabs — switch between All / YouTube / TikTok / Unused ideas",
    "Each idea card — tap to view, edit, or move to calendar",
    "Status badge — shows if idea is 'New', 'Scheduled', or 'Used'",
], 0.6, 2.0, 12, 4.0, size=15, dot_color=C_ACCENT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Videos Page
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.08, C_SKY)
add_text_box(slide, "Videos  —  Your Video Library", 0.6, 0.3, 12, 0.8, size=36, bold=True, color=C_WHITE)
add_text_box(slide, "Track every video you've recorded or posted", 0.65, 1.0, 12, 0.5, size=18, color=C_LIGHT)

add_text_box(slide, "What you can do:", 0.6, 1.6, 6, 0.4, size=16, bold=True, color=C_SKY)
add_bullet_box(slide, [
    "See all your videos in one place (YouTube + TikTok)",
    "Filter by status: All / Recording / Editing / Posted",
    "Tap a video card to add performance notes (views, likes, comments)",
    "Star icon — mark a video as a favourite",
    "Review button (magnifier icon) — open the AI Review Mode for that video",
    "Add Video button (+) — log a new video you've recorded",
], 0.6, 2.0, 12, 4.0, size=15, dot_color=C_SKY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Calendar Page
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.08, C_ROSE)
add_text_box(slide, "Calendar  —  Your Posting Schedule", 0.6, 0.3, 12, 0.8, size=36, bold=True, color=C_WHITE)
add_text_box(slide, "Plan which video gets posted on which day", 0.65, 1.0, 12, 0.5, size=18, color=C_LIGHT)

add_bullet_box(slide, [
    "Monthly calendar view — coloured dots show what's scheduled on each day",
    "Tap a date — see all posts planned for that day",
    "Drag and drop — move a post to a different day by holding and dragging it",
    "Add button (+) — schedule a new post on a specific date",
    "Each scheduled item — shows platform (YouTube/TikTok) and status",
    "Today is highlighted — easy to see where you are in the week",
    "Tap a scheduled item to change its status (Scheduled → Posted)",
], 0.6, 1.7, 12, 4.5, size=15, dot_color=C_ROSE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Leads (Client Pipeline)
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.08, C_ACCENT2)
add_text_box(slide, "Leads  —  Track Potential Clients", 0.6, 0.3, 12, 0.8, size=36, bold=True, color=C_WHITE)
add_text_box(slide, "Every person who reaches out or shows interest in your coaching", 0.65, 1.0, 12, 0.5, size=18, color=C_LIGHT)

add_text_box(slide, "The 7 stages (columns):", 0.6, 1.55, 12, 0.4, size=16, bold=True, color=C_ACCENT2)

stages = [
    ("New Lead",      C_SKY,    "Just reached out"),
    ("Contacted",     C_ACCENT, "You've replied"),
    ("Call Scheduled",C_YELLOW, "Call is booked"),
    ("Call Done",     RGBColor(0xF9,0x73,0x16), "Call happened"),
    ("Client",        C_ACCENT2,"Paying client"),
    ("Follow Up",     C_ROSE,   "Needs chasing"),
    ("Closed",        RGBColor(0x71,0x71,0x71),"No longer interested"),
]
for i, (name, col, desc) in enumerate(stages):
    x = 0.3 + i * 1.82
    add_rect(slide, x, 2.0, 1.6, 0.4, col)
    add_text_box(slide, name, x, 2.0, 1.6, 0.4, size=11, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, x, 2.5, 1.6, 0.6, size=11, color=C_LIGHT, align=PP_ALIGN.CENTER)

add_text_box(slide, "How to use:", 0.6, 3.3, 12, 0.4, size=16, bold=True, color=C_YELLOW)
add_bullet_box(slide, [
    "Add Lead (+) button — add a new person who contacted you",
    "Drag a card left or right to move them to the next stage",
    "Tap a lead card — edit their details, add notes, see activity history",
    "Content Attribution — record which video brought this person to you",
], 0.6, 3.75, 12, 2.5, size=15, dot_color=C_YELLOW)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Business Page
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.08, C_YELLOW)
add_text_box(slide, "Business  —  Your Growth Numbers", 0.6, 0.3, 12, 0.8, size=36, bold=True, color=C_WHITE)
add_text_box(slide, "See how your content is turning viewers into clients", 0.65, 1.0, 12, 0.5, size=18, color=C_LIGHT)

add_bullet_box(slide, [
    "Conversion Funnel — how many leads moved from New Lead all the way to Client",
    "Source Breakdown — which platform (TikTok, YouTube, WhatsApp) brings the most leads",
    "Content → Clients — which VIDEO category brings the most paying clients",
    "Views vs Inquiries chart — compare how content performance relates to enquiries",
    "Top Lead Videos — the specific videos that generated the most client interest",
], 0.6, 1.7, 12, 4.0, size=16, dot_color=C_YELLOW)

add_text_box(slide, "This page answers: \"Which videos are making me money?\"",
             0.6, 5.9, 12, 0.6, size=16, bold=True, color=C_ACCENT2)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Analytics
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.08, C_ACCENT)
add_text_box(slide, "Analytics  —  Performance Stats", 0.6, 0.3, 12, 0.8, size=36, bold=True, color=C_WHITE)
add_text_box(slide, "Charts and numbers showing how your content is performing", 0.65, 1.0, 12, 0.5, size=18, color=C_LIGHT)

add_bullet_box(slide, [
    "Top videos — ranked by views, likes, comments, saves",
    "Posting consistency chart — how regularly you've been posting each week",
    "Platform breakdown — YouTube vs TikTok performance side by side",
    "Content category performance — which topics (Discipline, Islamic Parenting, etc.) perform best",
    "YouTube category intelligence — AI-classified performance by content type",
    "Content → Clients section — connects analytics to lead generation",
], 0.6, 1.7, 12, 4.2, size=15, dot_color=C_ACCENT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — More Pages (Batch, Inbox, Strategist, Weekly Report)
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.08, C_SKY)
add_text_box(slide, "More Pages  —  Advanced Tools", 0.6, 0.3, 12, 0.8, size=36, bold=True, color=C_WHITE)
add_text_box(slide, "Tap 'More' in the bottom bar to access these", 0.65, 1.0, 12, 0.5, size=18, color=C_LIGHT)

pages = [
    ("Batch Plan",     C_ACCENT,  "Plan a whole week of content in one go — enter a theme, AI creates 7 ideas"),
    ("Inbox",          C_ACCENT2, "Log questions from your audience (WhatsApp/TikTok) — AI turns them into ideas"),
    ("Strategist",     C_YELLOW,  "AI content strategist — gives you a 7-day content roadmap and recommendations"),
    ("Weekly Report",  C_ROSE,    "AI summary of your week — wins, warnings, and what to focus on next week"),
    ("Connections",    C_SKY,     "Connect your YouTube channel — syncs your videos automatically"),
    ("Settings",       RGBColor(0x94,0xA3,0xB8), "Change your name, weekly goal, preferred platform, and AI coach tone"),
]
for i, (name, col, desc) in enumerate(pages):
    col_x = 0.6 + (i % 2) * 6.3
    row_y = 1.7 + (i // 2) * 1.7
    add_rect(slide, col_x, row_y, 1.5, 1.4, col)
    add_text_box(slide, name, col_x, row_y + 0.4, 1.5, 0.6, size=13, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, col_x + 1.6, row_y + 0.1, 4.5, 1.2, size=14, color=C_LIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — Settings
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.08, RGBColor(0x94,0xA3,0xB8))
add_text_box(slide, "Settings  —  Your Profile", 0.6, 0.3, 12, 0.8, size=36, bold=True, color=C_WHITE)
add_text_box(slide, "Personalise the app to match how you work", 0.65, 1.0, 12, 0.5, size=18, color=C_LIGHT)

add_bullet_box(slide, [
    "Display Name — your name shown in the dashboard greeting",
    "Weekly Goal — how many videos you aim to post per week (default: 5)",
    "Preferred Platform — YouTube or TikTok (affects which ideas are shown first)",
    "AI Coach Tone — choose between Motivational, Calm, or Direct coaching style",
    "Save button — tap after making any changes",
], 0.6, 1.7, 12, 4.0, size=16)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — Quick Reference Summary
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.08, C_ACCENT)
add_text_box(slide, "Quick Reference — Every Page", 0.6, 0.3, 12, 0.8, size=36, bold=True, color=C_WHITE)

pages_ref = [
    ("/today",          C_ACCENT2, "Daily post + posting checklist"),
    ("/dashboard",      C_YELLOW,  "Streak, stats, AI coach, momentum"),
    ("/ideas",          C_ACCENT,  "Content idea bank + voice capture"),
    ("/videos",         C_SKY,     "Video library + performance tracking"),
    ("/calendar",       C_ROSE,    "Schedule posts by day"),
    ("/leads",          C_ACCENT2, "Client pipeline (7 stages)"),
    ("/business",       C_YELLOW,  "Which content brings clients"),
    ("/analytics",      C_ACCENT,  "Charts and performance stats"),
    ("/batch",          C_SKY,     "Plan a full week of content"),
    ("/inbox",          C_ROSE,    "Audience Q&A → content ideas"),
    ("/strategist",     C_ACCENT2, "AI 7-day content roadmap"),
    ("/weekly-report",  C_YELLOW,  "AI weekly summary + score"),
]
for i, (route, col, desc) in enumerate(pages_ref):
    col_x = 0.4 + (i % 3) * 4.3
    row_y = 1.3 + (i // 3) * 1.4
    add_rect(slide, col_x, row_y, 1.7, 0.4, col)
    add_text_box(slide, route, col_x, row_y, 1.7, 0.4, size=12, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, col_x + 1.8, row_y + 0.05, 2.3, 0.4, size=13, color=C_WHITE)

# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
out = r"c:\Users\hp\OneDrive\Desktop\Guri_Dagan_Guide.pptx"
prs.save(out)
print(f"Saved: {out}")
